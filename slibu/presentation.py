from pptx import Presentation
from pptx.util import Centipoints, Cm, Emu, Inches, Mm, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import etree
from PIL import Image


def set_font_attr(run, key, value):
    run.font.__setattr__(key, value)

def set_font_color(run, key, value):
    rgb = RGBColor.from_string(value)
    run.font.color.rgb = rgb

def set_font_size(run, key, value):
    run.font.size = Pt(value)

def set_shape_bg_color(shape, key, value):
    rgb = RGBColor.from_string(value)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

def set_shape_bg_alpha(shape, key, value):
    fill_elt = shape._sp.spPr.solidFill
    alpha = round(value * 100000)
    for clr in ['srgbClr', 'schemeClr']:
        try:
            fill_elt.__getattribute__(clr).insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}alpha", val=str(alpha)))
        except AttributeError:
            pass

def set_unsupported(run, key, value):
    pass

def set_attr_on_pptx_object(obj, setter, key, options, tly):
    attr_key = key
    if type(setter) is tuple:
        setter, attr_key = setter
    value = options.get(key)
    tly_value = tly.get(key)
    if value is not None:
        setter(obj, attr_key, value)
    elif tly_value is not None:
        setter(obj, attr_key, tly_value)

def set_run_verbatim(run, key, value):
    if value == False:
        run.text = run.text.replace('\n',' ')

def set_tf_valign(tf, key, value):
    if value == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif value == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif value == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM


class PPTXPresentation:
    def __init__(self, template):
        # self.prs = Presentation()
        self.prs = Presentation(template)
        self.cphs = []
        self.tfs = []
        self.tlys = []

    def save(self, outfn):
        self.prs.save(outfn)

    def add_slide(self, layout="Title and Content", title=""):
        layout = self.get_layout_by_name(layout)
        self.slide = self.prs.slides.add_slide(layout)
        self.slide.shapes.title.text = title

        content_ph = self.find_content_placeholder()
        text_frame = None
        if content_ph is not None:
            text_frame = content_ph.text_frame
        self._push_cph(content_ph)
        self._push_tf(text_frame, True)
        self._push_tly(self.get_default_text_layout(), inherit=False)

    def add_paragraph(self):
        # if text frame is virgin, use first paragraph, otherwise create one
        tf, has_content = self.tf()
        if has_content:
            self.paragraph = tf.paragraphs[0]
        else:
            self.paragraph = tf.add_paragraph()
        tly = self.tly()
        if tly.get('indent_level') is not None:
            self.paragraph.level = tly['indent_level']
        if tly.get('show_bullet') is not None:
            if tly.get('list_order') == True:
                self.add_list_number()
            else:
                self.add_list_bullet()
        self._set_tf(tf, False)

    def add_text(self, text, **attrs):
        run = self.create_run()
        # TODO: make flag

        run.text = text

        tly = self.tly()

        attr_setters = {
                'verbatim': set_run_verbatim,
                'italic': set_font_attr,
                'bold': set_font_attr,
                'underline': set_font_attr,
                'fontname': (set_font_attr, 'name'),
                'font_size': set_font_size,
                'color': set_font_color,
                'border': set_unsupported,
                'bg_color': set_unsupported,
        }

        for key in [ 'verbatim', 'italic', 'bold', 'underline', 'fontname', 'font_size', 'color', 'border', 'bg_color' ]:
            setter = attr_setters.get(key)
            set_attr_on_pptx_object(run, setter, key, attrs, tly)


    def add_list_bullet(self):
        # TODO: find out what panose does
        # TODO: set marL on pPr
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buFont", typeface="Arial", panose="020B0604020202020204", pitchFamily="34", charset="0"))
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", char="•"))
        # self.paragraph._pPr.attrib['marL'] = "1200150" 
        self.paragraph._pPr.attrib['indent'] = "-285750"

    def add_list_number(self):
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum", type="arabicPeriod"))
        # TODO: fix this length
        self.paragraph._pPr.attrib['indent'] = "-285750"
        # self.paragraph._pPr.attrib['marL'] = "1200150" 
        # self.paragraph._pPr.attrib['indent'] = "-5000"

    def create_run(self):
        run = self.paragraph.add_run()
        return run

    def add_picture(self, src, left=None, top=None, width=None, height=None, **options):
        left = self.make_length('left', left)
        top = self.make_length('top', top)
        width = self.make_length('width', width)
        height = self.make_length('height', height)

        if left is None:
            left = self.cph().left
        if top is None:
            top = self.cph().top

        if width is None and height is None:
            # TODO: fix this so that in img tags, we use the actual dimensions, not scaling
            width = self.get_cph_width()
            height = self.get_cph_height()
            img = Image.open(src)
            actual_width = img.width
            actual_height = img.height
            if actual_width/width <= actual_height/height:
                width = None
            else:
                height = None

        picture = self.slide.shapes.add_picture(src,
                left, top, width, height)

        rotation = options.get('rotation')
        if rotation is not None:
            picture.rotation = rotation

    def add_table(self, rows, columns):
        cph = self.cph()
        # shape = self.slide.shapes.add_table(rows, columns, cph.left, cph.top, cph.width, cph.height)
        shape = self.slide.shapes.add_table(rows, columns, cph.left, cph.top, cph.width, 0)
        self.table = shape.table
        # set column widths
        print(dir(self.table))
        print("TBL width", self.cph().width)
        table_width = self.cph().width

        col_widths = self.tly().get('column_widths')
        if col_widths is not None:
            col_total = sum(col_widths)
            col_pcts = [ c/col_total for c in col_widths ]
            # print("COL PCT", col_pcts)
            for i,c in enumerate(col_pcts):
                # print(dir(self.table.columns))
                # print(self.table.columns[i])
                self.table.columns[i].width = round(c * table_width)

    def start_table_cell(self, row, column):
        cell = self.table.cell(row, column)
        tf = cell.text_frame
        self._push_tf(tf, True)
        self.add_paragraph()
    
    def end_table_cell(self):
        self._pop_tf()

    def start_box(self, **options):

        if self.cph() is not None:
            dim = self.get_box_dimensions_from_shape(self.cph())
        else:
            dim = 0, 0, self.prs.slide_width, self.prs.slide_height
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options, 'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        box = self.slide.shapes.add_textbox(*dim)

        text_frame = box.text_frame

        # inherit word wrap setting from parent
        text_frame.word_wrap = None

        attr_setters = {
                'bg_color': set_shape_bg_color,
                'bg_alpha': set_shape_bg_alpha,
                'valign': set_tf_valign,
        }

        # shape options
        for key in ['bg_color', 'bg_alpha']:
            setter = attr_setters.get(key)
            set_attr_on_pptx_object(box, setter, key, options, self.tly())
        # TODO: other things such as alignment, word wrap, etc

        for key in ['valign']:
            setter = attr_setters.get(key)
            set_attr_on_pptx_object(text_frame, setter, key, options, self.tly())

        self._push_cph(box)
        self._push_tf(text_frame, True)
        self._push_tly(options)

    def end_box(self):
        self._pop_tly()
        self._pop_tf()
        self._pop_cph()

    def drop_slide(self):
        self._pop_tly()
        self._pop_tf()
        self._pop_cph()
        rId = self.prs.slides._sldIdLst[-1].rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[-1]

    def set_slide_layout(self, layout_name):
        title = self.slide.shapes.title.text
        tly = self.tly()
        self.drop_slide()
        self.add_slide(layout=layout_name, title=title)

    def set_slide_background(self, img_path, **options):
        dim = (0, 0, None, None)
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options, 'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        img = self.slide.shapes.add_picture(img_path, *dim)

        # This moves it to the background
        # self.slide.shapes._spTree.remove(img._element)
        # self.slide.shapes._spTree.insert(2, img._element)
        cursor_sp = self.slide.shapes[0]._element
        cursor_sp.addprevious(img._element)

    def has_content_placeholder(self):
        return self.cph() is not None

    def push_text_layout(self, **options):
        self._push_tly(options, inherit=True)

    def pop_text_layout(self):
        self._pop_tly()

    def set_text_layout(self, **options):
        self.tlys[-1].update(options)

    def get_layout_by_name(self, name):
        layouts = [ l for l in self.prs.slide_layouts if l.name == name ]
        if layouts != []:
            return layouts[0]
        return None

    def get_default_text_layout(self):
        return {
            'verbatim': False,
        }

    def get_box_dimensions_from_shape(self, shape):
        return shape.left, shape.top, shape.width, shape.height

    def update_box_dimensions(self, dimension, *new_dimension):
        dim = list(dimension)
        for i in range(len(dim)):
            if new_dimension[i] is not None:
                dim[i] = new_dimension[i]
        return dim


    def find_placeholder(self, phf_type):
        for shape in self.slide.shapes:
            if shape.is_placeholder:
                # print("DEBUG:", "ph_idx", shape.placeholder_format.idx)
                # print("DEBUG:", "ph_type", shape.placeholder_format.type)
                # print("DEBUG:", "ph_name", shape.name)
                if shape.placeholder_format.type == phf_type:
                    return shape
        return None

    def find_content_placeholder(self):
        return self.find_placeholder(7)
        # BODY=2, OBJECT=7

    def set_title_box(self, options):
        title_box = self.slide.shapes.title
        dim = self.get_box_dimensions_from_shape(title_box)
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options,'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        title_box.left = dim[0]
        title_box.top = dim[1]
        title_box.width = dim[2]
        title_box.height = dim[3]

        bg_color = options.get('bg_color')
        if bg_color is not None:
            rgb = RGBColor.from_string(bg_color)
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = rgb
        # set margins on title textframe
        
        margin_props = ['margin_left', 'margin_right', 'margin_top', 'margin_bottom' ]

        for p in margin_props:
            m = self.make_length_from_options(options, p)
            if m is not None:
                title_box.text_frame.__setattr__(p, m)

        # raise title to top
        cursor_sp = self.slide.shapes[-1]._element
        cursor_sp.addnext(title_box._element)


    def _push_cph(self, cph):
        self.cphs.append(cph)

    def _push_tf(self, tf, has_content):
        self.tfs.append((tf, has_content))

    def _set_tf(self, tf, has_content):
        self.tfs[-1] = (tf, has_content)

    def _push_tly(self, tly, inherit=True):
        if inherit:
            d = self.tlys[-1].copy()
        else:
            d = {}
        d.update(tly)
        self.tlys.append(d)

    def _pop_cph(self):
        self.cphs.pop()

    def _pop_tly(self):
        self.tlys.pop()

    def _pop_tf(self):
        self.tfs.pop()

    def cph(self):
        return self.cphs[-1]

    def tf(self):
        return self.tfs[-1]

    def tly(self):
        return self.tlys[-1]

    def get_cph_width(self):
        cph = self.cph()
        if cph is None:
            return self.prs.slide_width
        return cph.width

    def get_cph_height(self):
        cph = self.cph()
        if cph is None:
            return self.prs.slide_height
        return cph.height


    def length_percent(self, key, value):
        if key in ['left', 'margin_left', 'margin_right', 'width']:
            return value * self.prs.slide_width / 100
        else:
            return value * self.prs.slide_height / 100


    def make_length(self,key,length_with_unit):
        unit_factories = {
            'cp': Centipoints,
            'cm': Cm,
            'emu': Emu,
            'in': Inches,
            'mm': Mm,
            'pt': Pt
        }
        if length_with_unit is None:
            return None
        length, unit = length_with_unit
        if unit == '%':
            return round(self.length_percent(key, length))
        else:
            return round(unit_factories.get(unit, Inches)(length))

    def make_length_from_options(self, options, key):
        return self.make_length(key, options.get(key))




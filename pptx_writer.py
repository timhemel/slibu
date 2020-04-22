from pptx import Presentation
from pptx.util import Centipoints, Cm, Emu, Inches, Mm, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from lxml import etree
from PIL import Image
from pygments import highlight
from pygments.lexers import get_lexer_by_name
from pygments.formatter import Formatter
from pptx.dml.color import RGBColor
import re

class PresentationPygmentsFormatter(Formatter):

    def __init__(self, presentation, **options):
        Formatter.__init__(self, **options)
        self.presentation = presentation

    def format(self, tokensource, outfile):

        def get_font_attributes(ttype):
            d = {
                'italic': None,
                'bold': None,
                'underline': None,
                'fontname': None,
                'color': None,
                'border': None,
                'bg_color': None,
            }

            style = self.style.styles[ttype].split()
            for attr in style:
                # bold: render text as bold
                if attr == 'bold':
                    d['bold'] = True
                elif attr == 'nobold':
                    d['bold'] = False
                # italic: render text italic
                elif attr == 'italic':
                    d['italic'] = True
                elif attr == 'noitalic':
                    d['italic'] = False
                # underline render text underlined
                # nounderline don’t render text underlined
                elif attr == 'underline':
                    d['underline'] = True
                elif attr == 'nounderline':
                    d['underline'] = False
                # bg: transparent background
                # bg:#000000 background color (black)
                elif attr == 'bg:':
                    # background not supported on text run, 
                    d['bg_color'] = None
                    pass
                elif re.match(r'^bg:#[0-9A-Za-z]{6}', attr):
                    # pptx does not support background fill for text runs,
                    # only for complete shapes.
                    # rgb = RGBColor.from_string(attr[4:])
                    # run.font.fill.back_color.rgb = rgb
                    d['bg_color'] = attr[4:]
                # border: no border
                elif re.match(r'^border', attr):
                    d['border'] = None
                # border:#ffffff border color (white)
                elif re.match(r'^border:#[0-9A-Za-z]{6}', attr):
                    d['border'] = attr[8:]
                # #ff0000 text color (red)
                elif re.match(r'^#[0-9A-Za-z]{6}', attr) is not None:
                    # rgb = RGBColor.from_string(attr[1:])
                    # run.font.color.rgb = rgb
                    d['color'] = attr[1:]
                # noinherit don’t inherit styles from supertoken

            return d
 
        elements = [ (text, get_font_attributes(ttype)) for ttype, text in tokensource ]

        if elements[-1][0].rstrip('\n') == '':
            del elements[-1]

        for text, style in elements:
            self.presentation.add_text(text, verbatim=True, **style)


def set_font_attr(run, key, value):
    run.font.__setattr__(key, value)

def set_font_color(run, key, value):
    rgb = RGBColor.from_string(value)
    run.font.color.rgb = rgb

def set_font_size(run, key, value):
    run.font.size = Pt(value)

def set_shape_bg_color(tf, key, value):
    rgb = RGBColor.from_string(value)
    tf.fill.solid()
    tf.fill.fore_color.rgb = rgb

def set_unsupported(run, key, value):
    pass

def set_attr_on_pptx_object(obj, setter, key, options, tly):
    attr_key = key
    if type(setter) is tuple:
        setter, attr_key = setter
    value = options.get(key)
    tly_value = tly.get(key)
    if value is not None:
        setter(obj, attr_key, value)
    elif tly_value is not None:
        setter(obj, attr_key, tly_value)

def set_run_verbatim(run, key, value):
    if value == False:
        run.text = run.text.replace('\n',' ')



class PPTXPresentation:
    def __init__(self):
        # self.prs = Presentation()
        self.prs = Presentation("reference.pptx")
        self.cphs = []
        self.tfs = []
        self.tlys = []

    def save(self, outfn):
        self.prs.save(outfn)

    def add_slide(self, layout="Title and Content", title=""):
        layout = self.get_layout_by_name(layout)
        self.slide = self.prs.slides.add_slide(layout)
        self.slide.shapes.title.text = title

        content_ph = self.find_content_placeholder()
        text_frame = None
        if content_ph is not None:
            text_frame = content_ph.text_frame
        self._push_cph(content_ph)
        self._push_tf(text_frame, True)
        self._push_tly(self.get_default_text_layout(), inherit=False)

    def add_paragraph(self):
        # if text frame is virgin, use first paragraph, otherwise create one
        tf, has_content = self.tf()
        if has_content:
            self.paragraph = tf.paragraphs[0]
        else:
            self.paragraph = tf.add_paragraph()
        tly = self.tly()
        if tly.get('indent_level') is not None:
            self.paragraph.level = tly['indent_level']
        if tly.get('show_bullet') is not None:
            if tly.get('list_order') == True:
                self.add_list_number()
            else:
                self.add_list_bullet()
        self._set_tf(tf, False)

    def add_text(self, text, **attrs):
        run = self.create_run()
        # TODO: make flag

        run.text = text

        tly = self.tly()

        attr_setters = {
                'verbatim': set_run_verbatim,
                'italic': set_font_attr,
                'bold': set_font_attr,
                'underline': set_font_attr,
                'fontname': (set_font_attr, 'name'),
                'font_size': set_font_size,
                'color': set_font_color,
                'border': set_unsupported,
                'bg_color': set_unsupported,
        }

        for key in [ 'verbatim', 'italic', 'bold', 'underline', 'fontname', 'font_size', 'color', 'border', 'bg_color' ]:
            setter = attr_setters.get(key)
            set_attr_on_pptx_object(run, setter, key, attrs, tly)


    def add_list_bullet(self):
        # TODO: find out what panose does
        # TODO: set marL on pPr
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buFont", typeface="Arial", panose="020B0604020202020204", pitchFamily="34", charset="0"))
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar", char="•"))
        # self.paragraph._pPr.attrib['marL'] = "1200150" 
        self.paragraph._pPr.attrib['indent'] = "-285750"

    def add_list_number(self):
        self.paragraph._pPr.insert(0, etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum", type="arabicPeriod"))
        # TODO: fix this length
        self.paragraph._pPr.attrib['indent'] = "-285750"
        # self.paragraph._pPr.attrib['marL'] = "1200150" 
        # self.paragraph._pPr.attrib['indent'] = "-5000"

    def create_run(self):
        run = self.paragraph.add_run()
        return run

    def add_picture(self, src, left=None, top=None, width=None, height=None):
        left = self.make_length('left', left)
        top = self.make_length('top', top)
        width = self.make_length('width', width)
        height = self.make_length('height', height)

        if left is None:
            left = self.cph().left
        if top is None:
            top = self.cph().top

        if width is None and height is None:
            width = self.cph().width
            height = self.cph().height
            img = Image.open(src)
            actual_width = img.width
            actual_height = img.height
            if actual_width/width <= actual_height/height:
                width = None
            else:
                height = None

        picture = self.slide.shapes.add_picture(src,
                left, top, width, height)

    def add_table(self, rows, columns):
        cph = self.cph()
        # shape = self.slide.shapes.add_table(rows, columns, cph.left, cph.top, cph.width, cph.height)
        shape = self.slide.shapes.add_table(rows, columns, cph.left, cph.top, cph.width, 0)
        self.table = shape.table

    def start_table_cell(self, row, column):
        cell = self.table.cell(row, column)
        tf = cell.text_frame
        self._push_tf(tf, True)
        self.add_paragraph()
    
    def end_table_cell(self):
        self._pop_tf()

    def start_box(self, **options):

        dim = self.get_box_dimensions_from_shape(self.cph())
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options, 'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        box = self.slide.shapes.add_textbox(*dim)

        text_frame = box.text_frame

        # inherit word wrap setting from parent
        text_frame.word_wrap = None

        attr_setters = {
                'bg_color': set_shape_bg_color,
        }

        # shape options
        for key in ['bg_color']:
            setter = attr_setters.get(key)
            set_attr_on_pptx_object(box, setter, key, options, self.tly())
        # TODO: other things such as alignment, word wrap, etc

        self._push_cph(box)
        self._push_tf(text_frame, True)
        self._push_tly(options)

    def end_box(self):
        self._pop_tly()
        self._pop_tf()
        self._pop_cph()

    def drop_slide(self):
        self._pop_tly()
        self._pop_tf()
        self._pop_cph()
        rId = self.prs.slides._sldIdLst[-1].rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[-1]

    def set_slide_layout(self, layout_name):
        title = self.slide.shapes.title.text
        tly = self.tly()
        self.drop_slide()
        self.add_slide(layout=layout_name, title=title)

    def set_slide_background(self, img_path, **options):
        dim = (0, 0, None, None)
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options, 'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        img = self.slide.shapes.add_picture(img_path, *dim)

        # This moves it to the background
        # self.slide.shapes._spTree.remove(img._element)
        # self.slide.shapes._spTree.insert(2, img._element)
        cursor_sp = self.slide.shapes[0]._element
        cursor_sp.addprevious(img._element)

    def has_content_placeholder(self):
        return self.cph() is not None

    def push_text_layout(self, **options):
        self._push_tly(options, inherit=True)

    def pop_text_layout(self):
        self._pop_tly()

    def set_text_layout(self, **options):
        self.tlys[-1].update(options)

    def get_layout_by_name(self, name):
        layouts = [ l for l in self.prs.slide_layouts if l.name == name ]
        if layouts != []:
            return layouts[0]
        return None

    def get_default_text_layout(self):
        return {
            'verbatim': False,
        }

    def get_box_dimensions_from_shape(self, shape):
        return shape.left, shape.top, shape.width, shape.height

    def update_box_dimensions(self, dimension, *new_dimension):
        dim = list(dimension)
        for i in range(len(dim)):
            if new_dimension[i] is not None:
                dim[i] = new_dimension[i]
        return dim


    def find_placeholder(self, phf_type):
        for shape in self.slide.shapes:
            if shape.is_placeholder:
                # print("DEBUG:", "ph_idx", shape.placeholder_format.idx)
                # print("DEBUG:", "ph_type", shape.placeholder_format.type)
                # print("DEBUG:", "ph_name", shape.name)
                if shape.placeholder_format.type == phf_type:
                    return shape
        return None

    def find_content_placeholder(self):
        return self.find_placeholder(7)
        # BODY=2, OBJECT=7

    def set_title_box(self, options):
        title_box = self.find_placeholder(1)
        dim = self.get_box_dimensions_from_shape(title_box)
        dim = self.update_box_dimensions(dim,
                self.make_length_from_options(options,'left'),
                self.make_length_from_options(options, 'top'),
                self.make_length_from_options(options, 'width'),
                self.make_length_from_options(options, 'height'))

        title_box.left = dim[0]
        title_box.top = dim[1]
        title_box.width = dim[2]
        title_box.height = dim[3]

    def _push_cph(self, cph):
        self.cphs.append(cph)

    def _push_tf(self, tf, has_content):
        self.tfs.append((tf, has_content))

    def _set_tf(self, tf, has_content):
        self.tfs[-1] = (tf, has_content)

    def _push_tly(self, tly, inherit=True):
        if inherit:
            d = self.tlys[-1].copy()
        else:
            d = {}
        d.update(tly)
        self.tlys.append(d)

    def _pop_cph(self):
        self.cphs.pop()

    def _pop_tly(self):
        self.tlys.pop()

    def _pop_tf(self):
        self.tfs.pop()

    def cph(self):
        return self.cphs[-1]

    def tf(self):
        return self.tfs[-1]

    def tly(self):
        return self.tlys[-1]

    def length_percent(self, key, value):
        if key == 'left' or key == 'width':
            return value * self.prs.slide_width / 100
        else:
            return value * self.prs.slide_height / 100


    def make_length(self,key,length_with_unit):
        unit_factories = {
            'cp': Centipoints,
            'cm': Cm,
            'emu': Emu,
            'in': Inches,
            'mm': Mm,
            'pt': Pt
        }
        if length_with_unit is None:
            return None
        length, unit = length_with_unit
        if unit == '%':
            return round(self.length_percent(key, length))
        else:
            return round(unit_factories.get(unit, Inches)(length))

    def make_length_from_options(self, options, key):
        return self.make_length(key, options.get(key))



class PPTXWriter:

    def __init__(self):
        self.presentation = PPTXPresentation()

    def save(self, outfn):
        self.presentation.save(outfn)

    def feed(self, ast):
        for e in ast:
            print("DEBUG:", "e=", e)
            try:
                h = self.handlers[e['type']]
                try:
                    h(self, e)
                except Exception as e:
                    print("EXCEPTION", e)
                    raise(e)
            except KeyError:
                self.handle_undefined(e)

    def handle_undefined(self, e):
        print("ERROR: No handler for", e['type'])

    def handle_heading(self, e):
        if e['level'] == 1:
            self.presentation.add_slide(layout="1_Section Header", title=e['children'][0]['text'])
        elif e['level'] == 2:
            self.presentation.add_slide(layout="Title and Content", title=e['children'][0]['text'])

    def handle_paragraph(self, e):
        if self.presentation.has_content_placeholder():
            self.presentation.add_paragraph()
            self.feed(e['children'])

    def handle_block_text(self, e):
        if self.presentation.has_content_placeholder():
            self.presentation.add_paragraph()
            self.feed(e['children'])

    def handle_text(self, e):
        self.presentation.add_text(e['text'])

    def handle_emphasis(self, e):
        text = "".join([ x.get('text') for x in e['children'] ])
        self.presentation.add_text(text, italic=True)

    def handle_strong(self, e):
        text = "".join([ x.get('text') for x in e['children'] ])
        self.presentation.add_text(text, bold=True)

    def handle_codespan(self, e):
        self.presentation.add_text(e['text'], verbatim=True, fontname='Courier')

    def handle_block_quote(self, e):
        # TODO: fix fontname
        self.presentation.push_text_layout(italic=True, fontname='Comic Sans MS')
        self.feed(e['children'])
        self.presentation.pop_text_layout()

    def handle_list(self, e):
        self.presentation.push_text_layout(list_order=e['ordered'])
        self.feed(e['children'])
        self.presentation.pop_text_layout()

    def handle_list_item(self, e):
        self.presentation.push_text_layout(indent_level=e['level'], show_bullet=True)
        self.feed(e['children'])
        self.presentation.pop_text_layout()

    def handle_image(self, e):
        self.presentation.add_picture(e['src'])

    def handle_block_code(self, e):
        if self.presentation.has_content_placeholder():
            self.presentation.push_text_layout(fontname='Courier')
            self.presentation.add_paragraph()
            text = e['text'].rstrip('\n')
            if e['info'] is None:
                self.presentation.add_text(text, verbatim=True)
            else:
                lexer = get_lexer_by_name(e['info'].strip())
                formatter = PresentationPygmentsFormatter(self.presentation)
                result = highlight(text, lexer, formatter)

            self.presentation.pop_text_layout()

    def handle_table(self, e):
        rows = len(e['children'][1]['children']) + 1
        cols = len(e['children'][0]['children'])
        self.presentation.add_table(rows, cols)
        self.row = 0
        self.feed(e['children'])

    def handle_table_head(self, e):
        self.col = 0
        for ce in e['children']:
            self.feed([ce])
            self.col += 1
        self.row += 1

    def handle_table_body(self, e):
        self.feed(e['children'])

    def handle_table_row(self, e):
        self.col = 0
        for ce in e['children']:
            self.feed([ce])
            self.col += 1
        self.row += 1

    def handle_table_cell(self, e):
        self.presentation.start_table_cell(self.row, self.col)
        self.feed(e['children'])
        self.presentation.end_table_cell()

    def handle_box(self, e):
        self.presentation.start_box(**e['options'])
        self.feed(e['children'])
        self.presentation.end_box()

    def handle_img(self, e):
        self.presentation.add_picture(e['src'], **e['options'])

    def handle_title(self, e):
        self.presentation.set_title_box(e['options'])

    def handle_textstyle(self, e):
        self.presentation.push_text_layout(**e['options'])
        self.feed(e['children'])
        self.presentation.pop_text_layout()

    def handle_setstyle(self, e):
        self.presentation.set_text_layout(**e['options'])

    def handle_linebreak(self, e):
        if self.presentation.has_content_placeholder():
            self.presentation.add_paragraph()

    def handle_slide(self, e):
        layout = e['options'].get('layout')
        if layout is not None:
            self.presentation.set_slide_layout(layout)

        background = e['options'].get('background')
        if background is not None:
            self.presentation.set_slide_background(background, **e['options'])

    # TODO:
    # source code from file
    handlers = {
        'heading': handle_heading,
        'paragraph': handle_paragraph,
        'text': handle_text,
        'emphasis': handle_emphasis,
        'strong': handle_strong,
        'image': handle_image,
        'list': handle_list,
        'list_item': handle_list_item,
        'block_text': handle_block_text,
        'img': handle_img,
        'box': handle_box,
        'table': handle_table,
        'table_head': handle_table_head,
        'table_body': handle_table_body,
        'table_row': handle_table_row,
        'table_cell': handle_table_cell,
        'block_code': handle_block_code,
        'codespan': handle_codespan,
        'block_quote': handle_block_quote,
        'title': handle_title,
        'textstyle': handle_textstyle,
        'setstyle': handle_setstyle,
        'linebreak': handle_linebreak,
        'slide': handle_slide,
    }



